"""Python PPTX generator — replaces generate_deck.js, no Node.js required."""

from io import BytesIO
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour palette ─────────────────────────────────────────────────────────────
BLUE        = "1B3C87"
PURPLE      = "7C3FA8"
WHITE       = "FFFFFF"
LIGHT_GREY  = "F4F5F7"
MID_GREY    = "8C93A0"
DARK        = "1A1A2E"
AMBER       = "D97706"
GREEN       = "059669"
FONT        = "Calibri"

# ── Helpers ───────────────────────────────────────────────────────────────────

def _s(val) -> str:
    if val is None:
        return ""
    if isinstance(val, list):
        return ", ".join(str(v) for v in val)
    return str(val)


def _rgb(hex6: str) -> RGBColor:
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _bg(slide, hex6: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex6)


def _rect(slide, x, y, w, h, fill=None, line=None, lw=0.75):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape


def _rrect(slide, x, y, w, h, fill=None, line=None, radius=0.08, lw=0.75):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    half_short = min(w, h) / 2
    adj_val = min(int(radius / half_short * 50000), 50000) if half_short > 0 else 5000
    sp_pr = shape._element.spPr
    prstGeom = sp_pr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        else:
            for gd in list(avLst.findall(qn("a:gd"))):
                avLst.remove(gd)
        gd = etree.SubElement(avLst, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {adj_val}")
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape


def _txt(slide, text, x, y, w, h, size=12, bold=False, italic=False,
         color=MID_GREY, align="left", valign="top", wrap=True,
         char_spacing=0, font=FONT):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left  = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top   = Emu(0)
    tf.margin_bottom = Emu(0)
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map.get(valign, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = _s(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    run.font.name = font
    if char_spacing:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(char_spacing * 100)))
    return txb


def _header(slide, title: str, fill=BLUE):
    _rect(slide, 0, 0, 10, 0.7, fill=fill)
    _txt(slide, "docebo", 0.4, 0.1, 2, 0.5, size=14, bold=True, color=WHITE)
    _txt(slide, title, 0, 0.1, 9.6, 0.5, size=14, bold=True, color=WHITE, align="right")


def _footer(slide, meta: dict, date_str: str, color="555588"):
    pass


# ── Slide builders ─────────────────────────────────────────────────────────────

def _slide1_title(prs, deck: dict, date_str: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, DARK)
    meta = deck.get("meta", {})
    ts   = deck.get("title_slide", {})

    _rect(slide, 0, 0, 10, 1.1, fill=PURPLE)
    _txt(slide, "docebo",           0.5, 0.25, 3,   0.6,  size=22, bold=True, color=WHITE)
    _txt(slide, date_str,           6.5, 0.3,  3,   0.5,  size=11, color=WHITE, align="right")
    _txt(slide, "CAMPAIGN KICKOFF", 0.6, 1.4,  9,   0.5,  size=12, bold=True, color=PURPLE, char_spacing=4)
    _txt(slide, _s(ts.get("product") or meta.get("product")),
                                    0.6, 1.9,  9,   1.1,  size=44, bold=True, color=WHITE)
    _txt(slide, _s(ts.get("goal")), 0.6, 3.05, 8,   0.55, size=18, color=MID_GREY)
    _txt(slide, date_str,           0.6, 3.63, 6,   0.22, size=12, color="444466")

    goal_str = _s(ts.get("goal", ""))
    stats = [
        {"label": "TARGET",   "value": goal_str.split(" in ")[0] or ""},
        {"label": "TIMELINE", "value": _s(ts.get("timeline", ""))},
        {"label": "MOTION",   "value": _s(ts.get("motion", "")).split(" —")[0]},
    ]
    for i, st in enumerate(stats):
        x = 0.6 + i * 3.1
        _rrect(slide, x, 3.85, 2.8, 1.2, fill="232347", line="3A3A6E", radius=0.08)
        _txt(slide, st["value"], x, 3.95, 2.8, 0.55, size=16, bold=True, color=WHITE, align="center")
        _txt(slide, st["label"], x, 4.52, 2.8, 0.3,  size=9,  color=MID_GREY, align="center", char_spacing=2)

    _txt(slide, "Never Stop Learning", 0.5, 5.25, 9, 0.3, size=10, italic=True, color="555588")


def _slide2_status(prs, deck: dict, date_str: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    meta = deck.get("meta", {})
    ss   = deck.get("status_slide", {})

    _header(slide, "CAMPAIGN STATUS", fill=BLUE)

    _txt(slide, "CAMPAIGN READINESS", 0.5, 0.95, 4.2, 0.35, size=9, color=MID_GREY, char_spacing=2)

    _rrect(slide, 0.5, 1.3, 4.2, 2.05, fill=LIGHT_GREY, line="E0E0E0", radius=0.08)
    score = _s(ss.get("readiness_score") if ss.get("readiness_score") is not None
               else meta.get("readiness_score", ""))
    _txt(slide, score,    0.7,  1.4,  1.2, 0.9, size=52, bold=True, color=BLUE)
    _txt(slide, "/ 100",  1.85, 1.75, 1.0, 0.5, size=18, color=MID_GREY)

    status_label = _s(ss.get("status_label") or meta.get("status_label", ""))
    _txt(slide, "🟡  " + status_label, 0.7, 2.35, 2.5, 0.35, size=13, bold=True, color=AMBER)

    scores = [
        {"label": "Structure", "val": _s(ss.get("structure_score") if ss.get("structure_score") is not None else meta.get("structure_score", "")) + " / 100", "color": GREEN},
        {"label": "Evidence",  "val": _s(ss.get("evidence_score")  if ss.get("evidence_score")  is not None else meta.get("evidence_score", ""))  + " / 100", "color": AMBER},
    ]
    for i, sc in enumerate(scores):
        y = 2.75 + i * 0.28
        _txt(slide, sc["label"], 0.7, y, 1.2, 0.27, size=11, color=MID_GREY)
        _txt(slide, sc["val"],   2.0, y, 1.5, 0.27, size=12, bold=True, color=sc["color"])

    _rrect(slide, 0.5, 3.45, 4.2, 0.85, fill="FEF3C7", line="FCD34D", radius=0.06)
    gaps = _s(ss.get("evidence_gaps") or meta.get("evidence_gaps", ""))
    _txt(slide, "Evidence gaps: " + gaps, 0.65, 3.52, 3.9, 0.65, size=10, color="92400E", wrap=True)

    _rect(slide, 0.5, 4.38, 4.2, 0.01, fill="D0D0D8")
    _txt(slide, "CAMPAIGN NAME", 0.5, 4.43, 4.2, 0.18, size=7.5, bold=True, color=PURPLE, char_spacing=1.5)
    _txt(slide, _s(ss.get("campaign_name") or meta.get("campaign_name", "")),
         0.5, 4.59, 4.2, 0.28, size=12, bold=True, color=DARK)
    _txt(slide, "CAMPAIGN CONCEPT", 0.5, 4.9, 4.2, 0.18, size=7.5, bold=True, color=MID_GREY, char_spacing=1.5)
    _txt(slide, _s(ss.get("campaign_concept") or meta.get("campaign_concept", "")),
         0.5, 5.06, 4.2, 0.22, size=8, color="666688", wrap=True)

    _txt(slide, "TARGET & MOTION", 5.3, 0.95, 4.2, 0.35, size=9, color=MID_GREY, char_spacing=2)
    icp_items = [
        {"label": "ICP",      "value": _s(ss.get("icp", ""))},
        {"label": "MOTION",   "value": _s(ss.get("motion", ""))},
        {"label": "GOAL",     "value": _s(ss.get("goal", ""))},
        {"label": "TIMELINE", "value": _s(ss.get("timeline", ""))},
    ]
    for i, item in enumerate(icp_items):
        y = 1.3 + i * 0.9
        _rrect(slide, 5.3, y, 4.2, 0.78, fill=WHITE, line="E8E8E8", radius=0.06)
        _txt(slide, item["label"], 5.5, y + 0.07, 1.2, 0.25, size=8,   bold=True, color=PURPLE, char_spacing=1)
        _txt(slide, item["value"], 5.5, y + 0.3,  3.8, 0.42, size=10.5, color=DARK, wrap=True)

    _footer(slide, meta, date_str, color=MID_GREY)


def _slide3_pillars(prs, deck: dict, date_str: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, DARK)
    meta = deck.get("meta", {})
    ps   = deck.get("pillars_slide", {})

    _header(slide, "MESSAGING PILLARS", fill=PURPLE)
    _txt(slide, _s(ps.get("positioning", "")), 0.5, 0.85, 9, 0.55, size=19, italic=True, color=WHITE)

    pillars = (ps.get("pillars") or [])[:3]
    for i, p in enumerate(pillars):
        x = 0.35 + i * 3.15
        _rrect(slide, x, 1.55, 2.95, 3.55, fill="232347", line="3A3A6E", radius=0.1)
        _txt(slide, str(i + 1).zfill(2), x + 0.2, 1.7,  0.7,  0.5,  size=22, bold=True, color=PURPLE)
        _txt(slide, _s(p.get("title")),   x + 0.2, 2.25, 2.55, 0.75, size=13, bold=True, color=WHITE, wrap=True)
        _txt(slide, _s(p.get("one_liner")), x + 0.2, 3.05, 2.55, 0.9, size=10.5, color=MID_GREY, wrap=True)

        proof_text   = _s(p.get("proof", ""))
        proof_needed = (_s(p.get("proof_status")) == "needed"
                        or "proof point needed" in proof_text.lower())
        pf = "3D2800" if proof_needed else "0D2B1A"
        pl = "A05A00" if proof_needed else "1A5C34"
        pc = AMBER    if proof_needed else "86EFAC"
        pt = ("⚠  PROOF POINT NEEDED" if proof_needed else "✓  " + proof_text)
        _rrect(slide, x + 0.2, 4.05, 2.55, 0.75, fill=pf, line=pl, radius=0.06)
        _txt(slide, pt, x + 0.2, 4.1, 2.55, 0.65, size=9, color=pc, align="center", valign="middle", wrap=True)

    _footer(slide, meta, date_str, color="555588")


def _slide4_assets(prs, deck: dict, date_str: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    meta = deck.get("meta", {})
    as_  = deck.get("asset_slide", {})

    _header(slide, "ASSET PLAN", fill=BLUE)
    _txt(slide, _s(as_.get("asset_count", "")) + " assets · graph-selected channels",
         0.5, 0.85, 9, 0.4, size=13, italic=True, color=MID_GREY)

    col_labels = ["ASSET", "FORMAT", "OWNER", "PURPOSE"]
    col_w      = [2.2,  2.0,  1.6,  3.7]
    col_x      = [0.4, 2.65, 4.7, 6.35]

    for i, label in enumerate(col_labels):
        _rect(slide, col_x[i], 1.35, col_w[i] - 0.05, 0.35, fill=BLUE)
        _txt(slide, label, col_x[i] + 0.1, 1.35, col_w[i] - 0.15, 0.35,
             size=9, bold=True, color=WHITE, char_spacing=1, valign="middle")

    assets = (as_.get("assets") or [])[:6]
    for ri, a in enumerate(assets):
        y  = 1.75 + ri * 0.62
        bg = WHITE if ri % 2 == 0 else LIGHT_GREY
        cells = [_s(a.get("name")), _s(a.get("format")), _s(a.get("owner")), _s(a.get("purpose", ""))]
        for ci, cell in enumerate(cells):
            _rect(slide, col_x[ci], y, col_w[ci] - 0.05, 0.55, fill=bg, line="E0E0E0")
            cell_color = DARK if ci == 0 else (PURPLE if ci == 2 else MID_GREY)
            _txt(slide, cell, col_x[ci] + 0.1, y + 0.05, col_w[ci] - 0.2, 0.45,
                 size=11 if ci == 0 else 10, bold=(ci == 0),
                 color=cell_color, wrap=True, valign="middle")

    if as_.get("evidence_note"):
        _rrect(slide, 0.4, 4.6, 9.2, 0.55, fill="EEF2FF", line="C7D2FE", radius=0.07)
        _txt(slide, _s(as_["evidence_note"]), 0.6, 4.65, 8.8, 0.45, size=9.5, color="3730A3", wrap=True)

    _footer(slide, meta, date_str, color=MID_GREY)


def _slide5_rollout(prs, deck: dict, date_str: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, DARK)
    meta = deck.get("meta", {})
    rs   = deck.get("rollout_slide", {})

    _header(slide, "PHASED ROLLOUT", fill=BLUE)

    phases = (rs.get("phases") or [])[:2]
    for pi, ph in enumerate(phases):
        x = 0.35 + pi * 4.9
        _rrect(slide, x, 0.8, 4.55, 4.4, fill="1C1C3A", line="3A3A6E", radius=0.1)
        name = (_s(ph.get("name")) or f"Phase {pi + 1}").upper()
        _txt(slide, name, x + 0.2, 0.9,  3,   0.3,  size=8,  bold=True, color=PURPLE, char_spacing=1)
        _txt(slide, _s(ph.get("days") or ph.get("weeks", "")),
                          x + 0.2, 1.22, 4.1, 0.35, size=17, bold=True, color=WHITE)
        _txt(slide, "↳ " + _s(ph.get("milestone", "")),
                          x + 0.2, 1.6,  4.1, 0.45, size=9.5, italic=True, color=MID_GREY, wrap=True)

        tasks = (ph.get("tasks") or [])[:3]
        for ti, task in enumerate(tasks):
            ty = 2.15 + ti * 0.62
            _rrect(slide, x + 0.2, ty, 4.1, 0.55, fill="28284A", line="3A3A6E", radius=0.06)
            _txt(slide, _s(task), x + 0.35, ty + 0.08, 3.75, 0.4, size=9, color=WHITE, wrap=True)

        if ph.get("checkpoint"):
            _rrect(slide, x + 0.2, 4.05, 4.1, 0.65, fill="1A2F1A", line="2D5A2D", radius=0.06)
            _txt(slide, "✓  CHECKPOINT", x + 0.35, 4.1,  2,    0.25, size=8, bold=True, color=GREEN, char_spacing=1)
            _txt(slide, _s(ph["checkpoint"]),  x + 0.35, 4.33, 3.75, 0.3,  size=8.5, color="86EFAC", wrap=True)

    _footer(slide, meta, date_str, color="555588")


def _slide6_metrics(prs, deck: dict, date_str: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    meta = deck.get("meta", {})
    ms   = deck.get("metrics_slide", {})

    _header(slide, "SUCCESS METRICS & GOVERNANCE", fill=PURPLE)

    _txt(slide, "SUCCESS METRICS", 0.5, 0.85, 4.3, 0.3, size=9, bold=True, color=MID_GREY, char_spacing=2)

    metrics = (ms.get("metrics") or [])[:4]
    for i, m in enumerate(metrics):
        y        = 1.2 + i * 0.78
        verified = m.get("verified") is True
        bg_c  = "F0FDF4" if verified else "FFFBEB"
        ln_c  = "86EFAC" if verified else "FCD34D"
        tx_c  = "166534" if verified else "92400E"
        _rrect(slide, 0.5, y, 4.3, 0.68, fill=bg_c, line=ln_c, radius=0.06)
        _txt(slide, _s(m.get("label", "")).upper(), 0.65, y + 0.05, 3.9, 0.2,  size=7.5, color=MID_GREY, char_spacing=1)
        _txt(slide, _s(m.get("value", "")),          0.65, y + 0.26, 3.9, 0.35, size=10, bold=verified, color=tx_c, wrap=True)

    _txt(slide, "HUMAN REVIEW CHECKPOINTS", 5.2, 0.85, 4.3, 0.3, size=9, bold=True, color=MID_GREY, char_spacing=2)

    checkpoints = (ms.get("checkpoints") or [])[:3]
    for i, cp in enumerate(checkpoints):
        y = 1.2 + i * 1.12
        _rrect(slide, 5.2, y, 4.3, 1.0, fill=LIGHT_GREY, line="E0E0E0", radius=0.08)
        _txt(slide, _s(cp.get("day", "")),    5.35, y + 0.08, 0.8, 0.3,  size=13, bold=True, color=BLUE)
        _txt(slide, _s(cp.get("teams", "")),  6.2,  y + 0.1,  3.1, 0.28, size=9.5, bold=True, color=PURPLE)
        _txt(slide, _s(cp.get("action", "")), 5.35, y + 0.42, 3.95, 0.5, size=9,  color=MID_GREY, wrap=True)

    _footer(slide, meta, date_str, color=MID_GREY)


# ── Public entry point ─────────────────────────────────────────────────────────

def generate_pptx(deck_content: dict) -> bytes:
    """Build the 6-slide deck and return raw PPTX bytes."""
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    meta     = deck_content.get("meta", {})
    date_str = meta.get("date") or datetime.now().strftime("%B %d, %Y")

    _slide1_title(prs, deck_content, date_str)
    _slide2_status(prs, deck_content, date_str)
    _slide3_pillars(prs, deck_content, date_str)
    _slide4_assets(prs, deck_content, date_str)
    _slide5_rollout(prs, deck_content, date_str)
    _slide6_metrics(prs, deck_content, date_str)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
